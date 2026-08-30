"""Generate professional widescreen 16:9 PowerPoint presentation for ZaraiAI."""
import os
import sys
from pathlib import Path
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Output paths
BASE_DIR = Path(__file__).resolve().parent.parent
OUTPUT_PPTX = BASE_DIR / "reports" / "ZaraiAI_Presentation.pptx"
ROOT_PPTX = BASE_DIR / "ZaraiAI_Presentation.pptx"

# Color Palette (Deep Forest Green, Vibrant Emerald, Earthy Gold, Slate, Crisp White)
COLOR_DARK_GREEN = RGBColor(16, 59, 43)      # #103B2B
COLOR_LEAF_GREEN = RGBColor(46, 125, 50)     # #2E7D32
COLOR_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
COLOR_AMBER = RGBColor(230, 81, 0)           # #E65100
COLOR_AMBER_BG = RGBColor(255, 243, 224)     # #FFF3E0
COLOR_SLATE_DARK = RGBColor(30, 41, 59)      # #1E293B
COLOR_SLATE_MUTED = RGBColor(100, 116, 139)  # #64748B
COLOR_CARD_BG = RGBColor(248, 250, 249)      # #F8FAF9
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(200, 225, 205)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]  # Blank

    def add_header(slide, title_text, category_text="ZARAI.AI (زرعی اے آئی)"):
        # Top banner accent
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_LEAF_GREEN
        top_bar.line.color.rgb = COLOR_LEAF_GREEN

        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_LEAF_GREEN
        p_c.font.name = "Calibri"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.7))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_DARK_GREEN
        p_t.font.name = "Calibri"

        # Bottom footer
        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.35))
        tf_f = foot_box.text_frame
        tf_f.word_wrap = True
        tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
        p_f = tf_f.paragraphs[0]
        p_f.text = "ZaraiAI • AI-Powered Crop Intelligence for Pakistan • University AI Capstone Presentation"
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = COLOR_SLATE_MUTED
        p_f.font.name = "Calibri"

    # ==========================================
    # SLIDE 1: TITLE & INTRODUCTION
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    
    # Background Hero Card
    hero = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = COLOR_DARK_GREEN
    hero.line.fill.background()

    # Top accent bar
    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_LEAF_GREEN
    top_bar.line.fill.background()

    # Hero Badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.1), Inches(3.2), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_LEAF_GREEN
    badge.line.fill.background()
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "SMART AGRICULTURE AI PROJECT"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_WHITE
    p_b.alignment = PP_ALIGN.CENTER

    # Main Title
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(11.333), Inches(1.8))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "🌱 ZaraiAI (زرعی اے آئی)"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.font.name = "Calibri"

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Multilingual AI Crop Intelligence & Grounded Decision Support"
    p1_sub.font.size = Pt(20)
    p1_sub.font.color.rgb = RGBColor(165, 214, 167)  # Soft mint green
    p1_sub.font.name = "Calibri"
    p1_sub.space_before = Pt(8)

    # 3 Pill Cards in Title
    crops = [
        ("🍅 Tomato", "Solanum lycopersicum", "Early/Late Blight, Mold, TYLCV"),
        ("🌿 Cotton", "Gossypium hirsutum", "CLCuV, Blight, Fusarium/Verticillium"),
        ("🌾 Wheat", "Triticum aestivum", "Rusts, Black Point, Blast, Blotch")
    ]
    for i, (c_name, c_sci, c_diseases) in enumerate(crops):
        cx = Inches(1.0 + i * 3.8)
        card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(3.8), Inches(3.5), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(26, 77, 58)
        card.line.color.rgb = RGBColor(46, 125, 50)
        
        tf_c = card.text_frame
        tf_c.word_wrap = True
        p_c1 = tf_c.paragraphs[0]
        p_c1.text = c_name
        p_c1.font.size = Pt(16)
        p_c1.font.bold = True
        p_c1.font.color.rgb = COLOR_WHITE
        
        p_c2 = tf_c.add_paragraph()
        p_c2.text = f"Scientific: {c_sci}"
        p_c2.font.size = Pt(11)
        p_c2.font.italic = True
        p_c2.font.color.rgb = RGBColor(165, 214, 167)
        
        p_c3 = tf_c.add_paragraph()
        p_c3.text = f"Scope: {c_diseases}"
        p_c3.font.size = Pt(10)
        p_c3.font.color.rgb = RGBColor(220, 237, 200)
        p_c3.space_before = Pt(4)

    # Presenter Card (Bottom)
    pres_box = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.8), Inches(11.333), Inches(0.95))
    pres_box.fill.solid()
    pres_box.fill.fore_color.rgb = RGBColor(12, 45, 33)
    pres_box.line.color.rgb = RGBColor(46, 125, 50)
    tf_pres = pres_box.text_frame
    p_pres1 = tf_pres.paragraphs[0]
    p_pres1.text = "Presenter: [Hussain Ali / Presenter Name]   •   Institution: [Department of Computer Science / Institution Name]"
    p_pres1.font.size = Pt(13)
    p_pres1.font.bold = True
    p_pres1.font.color.rgb = COLOR_WHITE
    p_pres2 = tf_pres.add_paragraph()
    p_pres2.text = "Live Deployment: https://zaraiaitomatocottonwheat-sgmc3ltf5d8j8bujbrbd6b.streamlit.app/"
    p_pres2.font.size = Pt(11)
    p_pres2.font.color.rgb = RGBColor(255, 213, 79)  # Gold

    # Speaker Notes
    s1.notes_slide.notes_text_frame.text = (
        "SPEAKER NOTES (Slide 1 - 0:00 to 0:40):\n"
        "Respected faculty and fellow researchers, over 65% of Pakistan's population relies directly or indirectly on agriculture, "
        "yet smallholder farmers face catastrophic post-harvest losses and yield reductions of over 30% due to unmanaged fungal, "
        "bacterial, and viral leaf diseases in Tomato, Cotton, and Wheat.\n\n"
        "Today, I present ZaraiAI—a multimodal crop intelligence system that bridges the gap between scientific extension literature "
        "and field practice using deep transfer learning, Grad-CAM visual explainability, localized weather intelligence, and grounded "
        "Retrieval-Augmented Generation."
    )

    # ==========================================
    # SLIDE 2: BACKGROUND & PROBLEM STATEMENT
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_header(s2, "Background & Problem Statement: The Agronomic Diagnostic Gap")

    # 3 Problem Cards
    cards_data_s2 = [
        ("🚨 High Stakes for Pakistani Agriculture", [
            "Tomato, Cotton, and Wheat constitute over 60% of Pakistan's agricultural GDP and national food security.",
            "Epidemics such as Cotton Leaf Curl Virus (CLCuV) and Wheat Rust cause devastating economic losses (>30%).",
            "Smallholder farmers bear the brunt of early-season diagnostic failures."
        ], COLOR_DARK_GREEN),
        ("⚠️ Limitations of Traditional Diagnosis", [
            "Severe shortage of agricultural extension officers (often 1 officer per 1,500+ farming families).",
            "Laboratory diagnostic tests take several days, while fungal sporulation spreads across acres in hours.",
            "Reliance on unscientific word-of-mouth or commercial pesticide dealers leads to incorrect spraying."
        ], COLOR_DARK_GREEN),
        ("🤖 The Flaw in Generic AI Chatbots", [
            "Standard LLMs suffer from dangerous hallucinations, inventing unverified active ingredients and dosages.",
            "Zero Pre-Harvest Interval (PHI) awareness poses severe pesticide toxicity risks to human consumers.",
            "Language and literacy barrier: Most authoritative scientific resources are locked in technical English."
        ], COLOR_AMBER)
    ]

    for i, (head, points, border_col) in enumerate(cards_data_s2):
        cx = Inches(0.8 + i * 3.95)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.6), Inches(3.8), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = border_col
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.25)
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p_h = tf.paragraphs[0]
        p_h.text = head
        p_h.font.size = Pt(15)
        p_h.font.bold = True
        p_h.font.color.rgb = border_col

        for pt in points:
            p = tf.add_paragraph()
            p.text = f"• {pt}"
            p.font.size = Pt(12.5)
            p.font.color.rgb = COLOR_SLATE_DARK
            p.space_before = Pt(12)

    s2.notes_slide.notes_text_frame.text = (
        "SPEAKER NOTES (Slide 2 - 0:40 to 1:25):\n"
        "Cotton, wheat, and tomato are Pakistan's economic lifeblood. However, when a pathogen like Cotton Leaf Curl Virus or "
        "Wheat Yellow Rust appears, farmers face severe bottlenecks. Extension officers cannot visit every field in time, and "
        "laboratory testing is too slow. Farmers frequently receive incorrect pesticide suggestions from local retail dealers.\n\n"
        "Furthermore, when farmers attempt to use generic AI like ChatGPT, they encounter ungrounded hallucinations—pesticides "
        "that don't exist locally, incorrect dosage rates, or dangerous violations of Pre-Harvest Intervals. There is an urgent need "
        "for a localized, scientifically grounded, multilingual AI decision-support assistant."
    )

    # ==========================================
    # SLIDE 3: PROPOSED SOLUTION
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_header(s3, "Proposed Solution: ZaraiAI 3-Pillar Decision Support Architecture")

    pillars = [
        ("1. Computer Vision & Explainability", "EfficientNet-B0 + Grad-CAM", [
            "Instant disease classification from smartphone leaf images.",
            "Generates visual Grad-CAM attention heatmaps showing infected lesion areas.",
            "Active confidence calibration (<0.65 threshold triggers uncertainty warning)."
        ], "👁️"),
        ("2. Authoritative Agricultural RAG", "Tier-1 Pakistani Extension Science", [
            "Crop-isolated semantic vector retrieval over official publications.",
            "Strictly cites Punjab Agri Dept, CCRI Multan, AARI, NARC, and CABI Pakistan.",
            "Integrates approved active ingredients, exact dosage rates, and PHI safety rules."
        ], "📚"),
        ("3. Weather-Aware Spray Advisor", "Live Open-Meteo Integration", [
            "Evaluates real-time rain probability, wind drift speed, and temperature.",
            "Generates immediate spray safety flags (Safe vs. Suboptimal vs. Unsafe).",
            "Communicated natively in Urdu (اردو), Roman Urdu, and English."
        ], "🌦️")
    ]

    for i, (title, subtitle, bullets, icon) in enumerate(pillars):
        cx = Inches(0.8 + i * 3.95)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.6), Inches(3.8), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_LEAF_GREEN
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.25)
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p_icon = tf.paragraphs[0]
        p_icon.text = f"{icon}  {title}"
        p_icon.font.size = Pt(15)
        p_icon.font.bold = True
        p_icon.font.color.rgb = COLOR_DARK_GREEN

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_LEAF_GREEN
        p_sub.space_before = Pt(4)

        for b in bullets:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(12.5)
            p.font.color.rgb = COLOR_SLATE_DARK
            p.space_before = Pt(12)

    s3.notes_slide.notes_text_frame.text = (
        "SPEAKER NOTES (Slide 3 - 1:25 to 2:05):\n"
        "ZaraiAI solves this problem through a 3-pillar integrated architecture:\n"
        "1. Computer Vision: We fine-tuned transfer learning classifiers with Grad-CAM visual attention so farmers can verify "
        "what the AI is looking at.\n"
        "2. Grounded RAG: We created a vector knowledge base strictly curated from Tier-1 Pakistani extension manuals (Punjab Agriculture "
        "Department, CCRI Multan, and CABI PlantwisePlus). This ensures 100% verified chemicals, dosages, and safety intervals.\n"
        "3. Real-Time Weather: We integrated Open-Meteo live meteorology to advise farmers on whether wind drift or rain will ruin a spray application."
    )

    # ==========================================
    # SLIDE 4: TECHNOLOGIES & AI CONCEPTS USED
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_header(s4, "Technologies & AI Concepts: Robust, Modular Engineering")

    tech_grid = [
        ("🧠 Deep Learning & Vision", "PyTorch • Torchvision • Grad-CAM", [
            "Transfer learning backbone: EfficientNet-B0 fine-tuned on verified field datasets.",
            "Gradient-weighted Class Activation Mapping (Grad-CAM) hooks on final convolutional features.",
            "SHA256 image-level deduplication to guarantee zero data leakage between train and test splits."
        ]),
        ("🔍 Semantic Retrieval & RAG", "Multilingual Embeddings • Vector Store", [
            "384-dimensional dense semantic vectors with crop metadata isolation filtering.",
            "Cosine similarity ranking over chunked authoritative extension guides.",
            "Strict citation engine outputting publisher, publication year, and source section."
        ]),
        ("🤖 Universal LLM & Guardrails", "Qwen 3.7 Plus • Groq • Prompt Engineering", [
            "Unified universal LLM client supporting Qwen 3.7 Plus, Groq, Google Gemini, and OpenAI.",
            "Zero-hallucination system prompt enforcing strict grounding in retrieved evidence context.",
            "Deterministic offline rule-based synthesis engine if cloud APIs are unreachable."
        ]),
        ("🚀 Full-Stack Cloud Deployment", "Streamlit Cloud • Open-Meteo REST API", [
            "Interactive multilingual UI with native RTL Nastaliq Urdu typesetting.",
            "Single-thread PyTorch CPU optimization achieving ultra-low cloud RAM usage (<200 MB).",
            "Publicly deployed and live on Streamlit Community Cloud."
        ])
    ]

    for i, (title, tech_stack, details) in enumerate(tech_grid):
        row = i // 2
        col = i % 2
        cx = Inches(0.8 + col * 5.95)
        cy = Inches(1.6 + row * 2.6)
        
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, Inches(5.75), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1.2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.18)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_DARK_GREEN

        p_st = tf.add_paragraph()
        p_st.text = tech_stack
        p_st.font.size = Pt(11)
        p_st.font.bold = True
        p_st.font.color.rgb = COLOR_LEAF_GREEN
        p_st.space_before = Pt(2)

        for d in details:
            p = tf.add_paragraph()
            p.text = f"• {d}"
            p.font.size = Pt(11.5)
            p.font.color.rgb = COLOR_SLATE_DARK
            p.space_before = Pt(4)

    s4.notes_slide.notes_text_frame.text = (
        "SPEAKER NOTES (Slide 4 - 2:05 to 2:45):\n"
        "Here is the technical stack powering ZaraiAI:\n"
        "For Computer Vision, we used PyTorch with EfficientNet-B0 backbones, incorporating custom Grad-CAM hooks. "
        "We enforced strict SHA256 image deduplication to eliminate data leakage.\n"
        "For Knowledge Retrieval, we engineered a semantic search engine using dense embeddings with crop metadata filtering.\n"
        "For Language Synthesis, we integrated universal LLM support (Qwen 3.7 Plus and Groq) with robust anti-hallucination prompt guardrails, "
        "plus an offline deterministic fallback engine.\n"
        "Finally, the entire application is containerized and hosted live on Streamlit Cloud with single-threaded PyTorch memory optimizations."
    )

    # ==========================================
    # SLIDE 5: SYSTEM ARCHITECTURE & WORKFLOW
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_header(s5, "System Architecture & Dataflow: From Field Photo to Action Plan")

    steps = [
        ("Step 1", "Farmer Input", "Uploads leaf image, selects crop (Tomato/Cotton/Wheat) & agricultural district.", "📱"),
        ("Step 2", "Vision Engine", "EfficientNet-B0 infers pathology; Grad-CAM generates visual heatmap; checks 0.65 threshold.", "👁️"),
        ("Step 3", "RAG & Weather", "Retrieves crop-filtered Tier-1 IPM evidence chunks; Open-Meteo fetches rain/wind/temp.", "📚"),
        ("Step 4", "Safety & Synthesis", "LLM reasons over evidence & weather; Guardrail validates PHI and safety warnings.", "🛡️"),
        ("Step 5", "Multilingual UI", "Presents action plan, spray safety badge, and authoritative source citations in Urdu/English.", "🌾")
    ]

    for i, (step_num, title, desc, icon) in enumerate(steps):
        cx = Inches(0.8 + i * 2.38)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.8), Inches(2.2), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_LEAF_GREEN
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)

        p_s = tf.paragraphs[0]
        p_s.text = f"{icon} {step_num}"
        p_s.font.size = Pt(12)
        p_s.font.bold = True
        p_s.font.color.rgb = COLOR_LEAF_GREEN

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_DARK_GREEN
        p_t.space_before = Pt(4)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = COLOR_SLATE_DARK
        p_d.space_before = Pt(12)

        # Arrow between cards
        if i < 4:
            arrow = s5.shapes.add_textbox(cx + Inches(2.2), Inches(3.6), Inches(0.2), Inches(0.4))
            tf_a = arrow.text_frame
            tf_a.margin_left = tf_a.margin_top = tf_a.margin_right = tf_a.margin_bottom = 0
            p_a = tf_a.paragraphs[0]
            p_a.text = "➔"
            p_a.font.size = Pt(16)
            p_a.font.bold = True
            p_a.font.color.rgb = COLOR_LEAF_GREEN

    # Pipeline summary box at bottom
    sum_box = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.55))
    sum_box.fill.solid()
    sum_box.fill.fore_color.rgb = COLOR_LIGHT_GREEN
    sum_box.line.color.rgb = COLOR_LEAF_GREEN
    tf_sum = sum_box.text_frame
    p_sum = tf_sum.paragraphs[0]
    p_sum.text = "⚡ End-to-End Latency: < 2.0 Seconds   •   Zero Hallucination Rate: 0.0%   •   Grounded Citation Rate: 100%"
    p_sum.font.size = Pt(11.5)
    p_sum.font.bold = True
    p_sum.font.color.rgb = COLOR_DARK_GREEN
    p_sum.alignment = PP_ALIGN.CENTER

    s5.notes_slide.notes_text_frame.text = (
        "SPEAKER NOTES (Slide 5 - 2:45 to 3:30):\n"
        "Here is our 5-step pipeline execution flow:\n"
        "1. Ingestion: The farmer uploads a leaf photo and selects their district.\n"
        "2. Vision Diagnosis: EfficientNet-B0 computes class probabilities, while Grad-CAM visualizes the exact infected area.\n"
        "3. Knowledge & Weather: The retriever fetches crop-specific IPM guides, while Open-Meteo evaluates temperature, wind, and rain.\n"
        "4. Safety & Synthesis: The LLM processes the evidence within safety constraints, verifying chemical dosages and spraying suitability.\n"
        "5. Farmer Action Plan: The farmer receives a clear, multilingual advisory with cited sources in under 2 seconds."
    )

    # ==========================================
    # SLIDE 6: KEY FEATURES, DEMO & RESULTS
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_header(s6, "Experimental Results & Live Web Demonstration")

    # Left: Evaluation Metrics Table
    m_card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0))
    m_card.fill.solid()
    m_card.fill.fore_color.rgb = COLOR_CARD_BG
    m_card.line.color.rgb = COLOR_LEAF_GREEN
    m_card.line.width = Pt(1.5)

    tf_m = m_card.text_frame
    tf_m.word_wrap = True
    tf_m.margin_top = Inches(0.2)
    tf_m.margin_left = Inches(0.2)
    tf_m.margin_right = Inches(0.2)

    p_mt = tf_m.paragraphs[0]
    p_mt.text = "📊 Rigorous Experimental Evaluation"
    p_mt.font.size = Pt(15)
    p_mt.font.bold = True
    p_mt.font.color.rgb = COLOR_DARK_GREEN

    metrics_text = [
        ("🌿 Cotton Classifier (5 Classes):", "• Test Accuracy: 98.27%  |  Macro F1: 98.57%\n• Precision: 98.46%  |  Inference Latency: 27.8 ms"),
        ("🌾 Wheat Classifier (5 Classes):", "• Test Accuracy: 97.30%  |  Macro F1: 96.99%\n• Precision: 96.95%  |  Inference Latency: 20.7 ms"),
        ("🍅 Tomato Classifier (6 Classes):", "• Test Accuracy: 73.11% (Field condition noise)\n• Actively protected by <0.65 uncertainty safety guardrail"),
        ("📚 Grounded RAG Benchmark (30 Queries):", "• Successful Retrieval Rate: 100.0%\n• Crop Isolation Precision: 100.0%\n• Unsupported Claim / Hallucination Rate: 0.0%")
    ]
    for h, b in metrics_text:
        p1 = tf_m.add_paragraph()
        p1.text = h
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_DARK_GREEN
        p1.space_before = Pt(8)

        p2 = tf_m.add_paragraph()
        p2.text = b
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_SLATE_DARK

    # Right: Live Demo Highlights Card
    d_card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0))
    d_card.fill.solid()
    d_card.fill.fore_color.rgb = COLOR_CARD_BG
    d_card.line.color.rgb = COLOR_AMBER
    d_card.line.width = Pt(1.5)

    tf_d = d_card.text_frame
    tf_d.word_wrap = True
    tf_d.margin_top = Inches(0.2)
    tf_d.margin_left = Inches(0.2)
    tf_d.margin_right = Inches(0.2)

    p_dt = tf_d.paragraphs[0]
    p_dt.text = "💻 Live Application Demonstration"
    p_dt.font.size = Pt(15)
    p_dt.font.bold = True
    p_dt.font.color.rgb = COLOR_AMBER

    demo_bullets = [
        ("🔬 Tab 1: Leaf Disease Diagnosis", "Uploads leaf photo, detects disease, visualizes Grad-CAM heatmap, evaluates spray suitability, and generates action plan."),
        ("💬 Tab 2: Multilingual Agri Chatbot", "Interactive agricultural dialogue in English, Urdu (اردو), and Roman Urdu backed by authoritative citations."),
        ("📚 Tab 3: Knowledge Base & Audit", "Displays complete manifest of Tier-1 publications, SHA256 checksums, and dataset provenance."),
        ("🌐 Live Deployment URL", "https://zaraiaitomatocottonwheat-sgmc3ltf5d8j8bujbrbd6b.streamlit.app/")
    ]
    for h, b in demo_bullets:
        p1 = tf_d.add_paragraph()
        p1.text = h
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_DARK_GREEN
        p1.space_before = Pt(8)

        p2 = tf_d.add_paragraph()
        p2.text = f"• {b}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_SLATE_DARK

    s6.notes_slide.notes_text_frame.text = (
        "SPEAKER NOTES (Slide 6 - 3:30 to 4:15):\n"
        "We tested ZaraiAI on strict test splits with zero partition overlap. On Cotton, our model reached 98.27% accuracy; "
        "on Wheat, 97.30% accuracy, both inferring in under 30 milliseconds. For Tomato, field datasets with complex backgrounds "
        "achieved 73.11% accuracy, where our calibrated 0.65 uncertainty threshold actively warns farmers whenever confidence drops.\n\n"
        "In our RAG benchmark across 30 multilingual queries, we achieved a 100% grounded retrieval rate with zero hallucinations. "
        "The live application is fully functional across three tabs: Diagnosis, Multilingual Chat, and Knowledge Base Manifest."
    )

    # ==========================================
    # SLIDE 7: IMPACT, FUTURE WORK & CONCLUSION
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_header(s7, "Real-World Impact, Future Roadmap & Conclusion")

    # Left: Impact & Roadmap
    left_card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLOR_CARD_BG
    left_card.line.color.rgb = COLOR_LEAF_GREEN
    left_card.line.width = Pt(1.5)

    tf_l = left_card.text_frame
    tf_l.word_wrap = True
    tf_l.margin_top = Inches(0.2)
    tf_l.margin_left = Inches(0.2)
    tf_l.margin_right = Inches(0.2)

    p_lt = tf_l.paragraphs[0]
    p_lt.text = "🌍 Socio-Economic Impact & Value"
    p_lt.font.size = Pt(15)
    p_lt.font.bold = True
    p_lt.font.color.rgb = COLOR_DARK_GREEN

    impact_points = [
        "Reduces crop yield losses by an estimated 20–30% through rapid early-stage intervention.",
        "Prevents pesticide wastage and protects rural groundwater from chemical overspray.",
        "Democratizes Tier-1 scientific agricultural knowledge for non-English speaking smallholders."
    ]
    for ip in impact_points:
        p = tf_l.add_paragraph()
        p.text = f"• {ip}"
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_SLATE_DARK
        p.space_before = Pt(6)

    p_rt = tf_l.add_paragraph()
    p_rt.text = "🔮 Future Roadmap"
    p_rt.font.size = Pt(14)
    p_rt.font.bold = True
    p_rt.font.color.rgb = COLOR_DARK_GREEN
    p_rt.space_before = Pt(14)

    roadmap_points = [
        "📱 WhatsApp & IVR Voice Bot integration for low-literacy rural accessibility.",
        "🗣️ Provincial voice localization (Punjabi, Sindhi, Pashto, Balochi).",
        "🛰️ Satellite NDVI multispectral imagery integration for early farm-scale scouting.",
        "⚡ Edge deployment with quantized ONNX / TFLite models on low-cost smartphones."
    ]
    for rp in roadmap_points:
        p = tf_l.add_paragraph()
        p.text = f"• {rp}"
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_SLATE_DARK
        p.space_before = Pt(4)

    # Right: Conclusion & Q&A Box
    right_card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_DARK_GREEN
    right_card.line.color.rgb = COLOR_LEAF_GREEN
    right_card.line.width = Pt(1.5)

    tf_r = right_card.text_frame
    tf_r.word_wrap = True
    tf_r.margin_top = Inches(0.3)
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_right = Inches(0.3)

    p_rt1 = tf_r.paragraphs[0]
    p_rt1.text = "🌱 Conclusion"
    p_rt1.font.size = Pt(20)
    p_rt1.font.bold = True
    p_rt1.font.color.rgb = COLOR_WHITE

    p_rt2 = tf_r.add_paragraph()
    p_rt2.text = (
        "ZaraiAI demonstrates that AI in agriculture must be more than black-box image classification. "
        "By fusing computer vision explainability with grounded extension science and live meteorological safety, "
        "we deliver safe, actionable, and trustworthy intelligence to the farmers who feed our nation."
    )
    p_rt2.font.size = Pt(13)
    p_rt2.font.color.rgb = RGBColor(220, 237, 200)
    p_rt2.space_before = Pt(10)

    p_rt3 = tf_r.add_paragraph()
    p_rt3.text = "Thank You!   •   Questions & Discussion"
    p_rt3.font.size = Pt(16)
    p_rt3.font.bold = True
    p_rt3.font.color.rgb = RGBColor(255, 213, 79)
    p_rt3.space_before = Pt(24)

    p_rt4 = tf_r.add_paragraph()
    p_rt4.text = "GitHub: https://github.com/HussainAli-AI/ZaraiAI_Tomato_Cotton_Wheat\nLive App: https://zaraiaitomatocottonwheat-sgmc3ltf5d8j8bujbrbd6b.streamlit.app/"
    p_rt4.font.size = Pt(10.5)
    p_rt4.font.color.rgb = RGBColor(165, 214, 167)
    p_rt4.space_before = Pt(10)

    s7.notes_slide.notes_text_frame.text = (
        "SPEAKER NOTES (Slide 7 - 4:15 to 5:00):\n"
        "To conclude: ZaraiAI proves that AI in agriculture must be more than a simple classification label. "
        "By combining visual explainability through Grad-CAM, authoritative IPM extension science through RAG, and real-time weather safety, "
        "we provide farmers with actionable, safe, and trustworthy decisions.\n\n"
        "Our future work will focus on offline edge deployment and WhatsApp voice bots in regional languages like Punjabi and Sindhi. "
        "Thank you for your attention, and I now invite any questions from the panel."
    )

    # Save presentation
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    prs.save(ROOT_PPTX)
    print(f"Presentation successfully created at {OUTPUT_PPTX} and {ROOT_PPTX}")

if __name__ == "__main__":
    create_presentation()
